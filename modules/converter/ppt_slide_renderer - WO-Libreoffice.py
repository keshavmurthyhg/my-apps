import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# -----------------------------------
# Skip unwanted slides
# -----------------------------------
def should_skip_slide(slide):
    texts = []

    for shape in slide.shapes:
        try:
            if hasattr(shape, "text"):
                txt = shape.text.strip().lower()
                if txt:
                    texts.append(txt)
        except:
            pass

    combined_text = " ".join(texts)

    skip_keywords = [
        "thank you",
        "questions",
        "ppt slides",
        "backup slides"
    ]

    for keyword in skip_keywords:
        if keyword in combined_text:
            return True

    return False


# -----------------------------------
# Detect background image
# -----------------------------------
def is_background_picture(shape, slide_width, slide_height):
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return False

        # very large image covering full slide
        if (
            shape.width >= slide_width * 0.90 and
            shape.height >= slide_height * 0.90 and
            shape.left <= slide_width * 0.05 and
            shape.top <= slide_height * 0.05
        ):
            return True

        return False

    except Exception:
        return False


# -----------------------------------
# Extract images directly from PPT
# -----------------------------------
def extract_images_from_slide(
    slide,
    slide_index,
    slide_width,
    slide_height,
    output_dir
):
    extracted_images = []
    image_counter = 1

    for shape in slide.shapes:
        try:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            # skip full-slide background images
            if is_background_picture(
                shape,
                slide_width,
                slide_height
            ):
                print(
                    f"Skipping background image in slide {slide_index}"
                )
                continue

            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext

            image_path = os.path.join(
                output_dir,
                f"slide_{slide_index}_{image_counter}.{image_ext}"
            )

            with open(image_path, "wb") as f:
                f.write(image_bytes)

            extracted_images.append(image_path)
            image_counter += 1

            print(
                f"Extracted image from slide {slide_index}"
            )

        except Exception as e:
            print(
                f"Image extraction failed "
                f"for slide {slide_index}: {e}"
            )

    return extracted_images


# -----------------------------------
# Create placeholder image when no images found
# -----------------------------------
def create_placeholder_image(
    slide_index,
    output_dir
):
    try:
        from PIL import Image, ImageDraw

        img = Image.new(
            "RGB",
            (1280, 720),
            color="white"
        )

        draw = ImageDraw.Draw(img)

        text = (
            f"Slide {slide_index}\n"
            f"(No embedded image found)"
        )

        draw.text(
            (100, 300),
            text,
            fill="black"
        )

        placeholder_path = os.path.join(
            output_dir,
            f"slide_{slide_index}_placeholder.png"
        )

        img.save(placeholder_path)

        return placeholder_path

    except Exception as e:
        print(
            f"Placeholder creation failed: {e}"
        )
        return None


# -----------------------------------
# Main function
# -----------------------------------
def render_ppt_slides_to_images(ppt_path):
    """
    Extracts images directly from PPT slides
    without LibreOffice dependency.
    """

    try:
        prs = Presentation(ppt_path)

        temp_dir = tempfile.mkdtemp()

        final_images = []

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_index, slide in enumerate(
            prs.slides,
            start=1
        ):

            print(
                f"Processing slide {slide_index}"
            )

            # skip unwanted slides
            if should_skip_slide(slide):
                print(
                    f"Skipping unwanted slide "
                    f"{slide_index}"
                )
                continue

            slide_images = extract_images_from_slide(
                slide=slide,
                slide_index=slide_index,
                slide_width=slide_width,
                slide_height=slide_height,
                output_dir=temp_dir
            )

            if slide_images:
                final_images.extend(slide_images)

            else:
                print(
                    f"No images found in "
                    f"slide {slide_index}"
                )

                placeholder = create_placeholder_image(
                    slide_index,
                    temp_dir
                )

                if placeholder:
                    final_images.append(
                        placeholder
                    )

        print(
            f"Total extracted images: "
            f"{len(final_images)}"
        )

        return final_images

    except Exception as e:
        print(
            f"PPT render failed: {e}"
        )
        return []