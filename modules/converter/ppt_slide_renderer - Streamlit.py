import os
import copy
import tempfile
import subprocess

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pdf2image import convert_from_path


# -----------------------------------
# Skip unwanted slides
# -----------------------------------
def should_skip_slide(slide):
    texts = []

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            txt = shape.text.strip().lower()
            if txt:
                texts.append(txt)

    combined_text = " ".join(texts)

    if "thank you" in combined_text:
        return True

    if "ppt slides" in combined_text:
        return True

    if "questions" in combined_text:
        return True

    return False


# -----------------------------------
# Detect background images
# -----------------------------------
def is_background_picture(shape, slide_width, slide_height):
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return False

        if (
            shape.width >= slide_width * 0.90 and
            shape.height >= slide_height * 0.90 and
            shape.left <= slide_width * 0.05 and
            shape.top <= slide_height * 0.05
        ):
            return True

        return False

    except:
        return False


# -----------------------------------
# Create clean ppt
# -----------------------------------
def create_clean_ppt(ppt_path):
    source_prs = Presentation(ppt_path)

    clean_prs = Presentation()
    clean_prs.slide_width = source_prs.slide_width
    clean_prs.slide_height = source_prs.slide_height

    blank_layout = clean_prs.slide_layouts[6]

    for slide in source_prs.slides:

        if should_skip_slide(slide):
            print("Skipping unwanted slide")
            continue

        new_slide = clean_prs.slides.add_slide(blank_layout)

        # ----------------------------
        # Step 1: Add screenshots/images
        # ----------------------------
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                    if is_background_picture(
                        shape,
                        source_prs.slide_width,
                        source_prs.slide_height
                    ):
                        print("Skipping background image")
                        continue

                    image = shape.image
                    image_bytes = image.blob

                    temp_img = tempfile.NamedTemporaryFile(
                        delete=False,
                        suffix="." + image.ext
                    )
                    temp_img.write(image_bytes)
                    temp_img.close()

                    new_slide.shapes.add_picture(
                        temp_img.name,
                        shape.left,
                        shape.top,
                        shape.width,
                        shape.height
                    )

            except Exception as e:
                print(f"Picture extraction failed: {e}")

        # ----------------------------
        # Step 2: Add annotations/shapes
        # ----------------------------
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue

                if shape.is_placeholder:
                    continue

                el = shape.element
                new_el = copy.deepcopy(el)

                new_slide.shapes._spTree.insert_element_before(
                    new_el,
                    'p:extLst'
                )

            except Exception as e:
                print(f"Annotation copy failed: {e}")

    temp_dir = tempfile.mkdtemp()

    clean_ppt_path = os.path.join(
        temp_dir,
        "clean_ppt.pptx"
    )

    clean_prs.save(clean_ppt_path)

    return clean_ppt_path, temp_dir


# -----------------------------------
# Convert clean ppt → images
# -----------------------------------
def render_ppt_slides_to_images(ppt_path):
    try:
        clean_ppt_path, temp_dir = create_clean_ppt(ppt_path)

        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            clean_ppt_path,
            "--outdir",
            temp_dir
        ], check=True)

        pdf_files = [
            os.path.join(temp_dir, f)
            for f in os.listdir(temp_dir)
            if f.endswith(".pdf")
        ]

        if not pdf_files:
            raise Exception("PDF conversion failed")

        pdf_path = pdf_files[0]

        pages = convert_from_path(
            pdf_path,
            dpi=200
        )

        final_images = []

        for i, page in enumerate(pages):
            img_path = os.path.join(
                temp_dir,
                f"slide_{i+1}.png"
            )

            page.save(img_path, "PNG")
            final_images.append(img_path)

        return final_images

    except Exception as e:
        print(f"PPT render failed: {e}")
        return []
